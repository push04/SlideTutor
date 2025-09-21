# SlideTutor — Improved version with EasyOCR (no Tesseract binary required)
# UI/UX refreshed: single-page (no sidebar), top tabs, modern CSS, subtle animations
# Save this file and run with: streamlit run slidetutor_improved_easyocr_ui_cleaned.py
from __future__ import annotations
import os
import io
import sys
import time
import json
import math
import base64
import tempfile
import traceback
import sqlite3
import logging
from typing import List, Dict, Tuple, Optional, Any

# third-party requirements (streamlit is required for UI)
try:
    import streamlit as st
except Exception as e:
    raise RuntimeError("streamlit must be installed. pip install streamlit") from e

# Optional libraries (single check location)
_HAS_FAISS = True
try:
    import faiss
except Exception:
    faiss = None
    _HAS_FAISS = False

_HAS_PPTX = True
try:
    from pptx import Presentation
except Exception:
    Presentation = None
    _HAS_PPTX = False

_HAS_PYMUPDF = True
try:
    import fitz
except Exception:
    fitz = None
    _HAS_PYMUPDF = False

_HAS_EASYOCR = True
try:
    import easyocr
    from PIL import Image
except Exception:
    easyocr = None
    Image = None
    _HAS_EASYOCR = False

_HAS_SENTENCE_TRANSFORMERS = True
try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None
    _HAS_SENTENCE_TRANSFORMERS = False

try:
    import numpy as np
except Exception:
    raise RuntimeError("numpy is required. pip install numpy")

_HAS_SKLEARN = True
try:
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:
    cosine_similarity = None
    _HAS_SKLEARN = False

_HAS_GTTS = True
try:
    from gtts import gTTS
except Exception:
    gTTS = None
    _HAS_GTTS = False

import requests

# logging
logger = logging.getLogger(__name__)
if not logger.handlers:
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(name)s: %(message)s"))
    logger.addHandler(handler)
logger.setLevel(logging.INFO)

# ------------------------------
# Configuration
# ------------------------------
DEFAULT_OPENROUTER_KEY = (st.secrets.get("OPENROUTER_API_KEY") if hasattr(st, 'secrets') else None) or os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_API_URL = os.getenv("OPENROUTER_API_URL", "https://openrouter.ai/api/v1/chat/completions")
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "all-MiniLM-L6-v2")
TOP_K = int(os.getenv("TOP_K", "5"))
DB_PATH = os.getenv("SLIDETUTOR_DB", "slidetutor.sqlite3")

APP_TITLE = "SlideTutor — Student PPT/PDF Tutor (Improved)"
APP_SUBTITLE = "Upload slides and get deep multi-level lessons, quizzes, flashcards and spaced repetition."

# ------------------------------
# Helpers
# ------------------------------

def log(*args, **kwargs):
    print(*args, **kwargs, file=sys.stderr)


def safe_json_loads(s: str) -> Any:
    try:
        return json.loads(s)
    except Exception:
        return None


def ensure_dir(path: str):
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)

# ------------------------------
# DB
# ------------------------------

def init_db(path: str = DB_PATH):
    conn = sqlite3.connect(path, check_same_thread=False)
    c = conn.cursor()
    c.execute(
        """
        CREATE TABLE IF NOT EXISTS uploads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT,
            uploaded_at INTEGER,
            meta TEXT
        )
        """
    )
    c.execute(
        """
        CREATE TABLE IF NOT EXISTS flashcards (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            upload_id INTEGER,
            question TEXT,
            answer TEXT,
            easiness REAL DEFAULT 2.5,
            interval INTEGER DEFAULT 1,
            repetitions INTEGER DEFAULT 0,
            next_review INTEGER
        )
        """
    )
    c.execute(
        """
        CREATE TABLE IF NOT EXISTS quizzes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            upload_id INTEGER,
            question TEXT,
            options TEXT,
            correct_index INTEGER,
            created_at INTEGER
        )
        """
    )
    conn.commit()
    return conn

_db_conn = init_db()


def upload_db_id(upload: Dict) -> Optional[int]:
    """
    Return DB uploads.id for this upload dict if available.
    Prefer explicit 'db_id' key, otherwise attempt a DB lookup by filename.
    Returns None if not resolvable.
    """
    if not upload:
        return None

    # Prefer explicit db_id if present and valid
    dbid = upload.get("db_id")
    if dbid is not None:
        try:
            return int(dbid)
        except Exception:
            # fallthrough to DB lookup
            logger.debug("upload_db_id: invalid db_id in upload object: %s", dbid)

    # Defensive DB lookup using cursor helper if available
    try:
        cur = _get_db_cursor()
        if cur is None:
            return None
        cur.execute(
            "SELECT id FROM uploads WHERE filename = ? ORDER BY uploaded_at DESC LIMIT 1",
            (upload.get("filename"),),
        )
        row = cur.fetchone()
        if row:
            return int(row[0])
    except Exception as e:
        logger.debug("upload_db_id lookup failed: %s", e)
    return None

# ------------------------------
# Vector index & embeddings
# ------------------------------

def _ensure_2d(arr: np.ndarray) -> np.ndarray:
    arr = np.asarray(arr)
    if arr.ndim == 1:
        return arr.reshape(1, -1)
    return arr


def _safe_normalize(a: np.ndarray, axis: int = 1, eps: float = 1e-9) -> np.ndarray:
    a = np.asarray(a, dtype=np.float32, order="C")
    norm = np.linalg.norm(a, axis=axis, keepdims=True)
    norm[norm == 0] = eps
    return a / norm


class VectorIndexFallback:
    def __init__(self, embeddings: Optional[np.ndarray], texts: Optional[List[str]] = None) -> None:
        texts = texts or []
        if embeddings is None or getattr(embeddings, "size", 0) == 0:
            self.embeddings = np.zeros((0, 0), dtype=np.float32)
            self._normed = None
        else:
            emb = np.asarray(embeddings, dtype=np.float32, order="C")
            if emb.ndim == 1:
                emb = emb.reshape(1, -1)
            self.embeddings = emb
            self._normed = _safe_normalize(self.embeddings, axis=1)
        self.texts = list(texts)

    def _empty_result(self, n_queries: int, k: int) -> Tuple[np.ndarray, np.ndarray]:
        dists = np.full((n_queries, k), np.finfo(np.float32).max, dtype=np.float32)
        idxs = np.full((n_queries, k), -1, dtype=np.int64)
        return dists, idxs

    def search(self, q_emb: np.ndarray, k: int = 5) -> Tuple[np.ndarray, np.ndarray]:
        q = np.asarray(q_emb, dtype=np.float32, order="C")
        q = _ensure_2d(q)
        n_q, dim_q = q.shape

        if self.embeddings.size == 0:
            return self._empty_result(n_q, k)

        if self.embeddings.shape[1] != dim_q:
            logger.error("Query embedding dimension (%d) does not match index dimension (%d).", dim_q, self.embeddings.shape[1])
            return self._empty_result(n_q, k)

        normed_q = _safe_normalize(q, axis=1)

        if _HAS_SKLEARN and cosine_similarity is not None and self._normed is not None:
            try:
                sims = cosine_similarity(normed_q, self._normed)
            except Exception as ex:
                logger.warning("sklearn cosine_similarity failed, falling back to numpy: %s", ex)
                sims = normed_q @ self._normed.T
        else:
            sims = normed_q @ self._normed.T

        n_db = sims.shape[1]
        k_eff = min(max(1, int(k)), n_db)
        idxs_part = np.argpartition(-sims, kth=k_eff - 1, axis=1)[:, :k_eff]
        row_indices = np.arange(n_q)[:, None]
        idxs_sorted = idxs_part[np.arange(n_q)[:, None], np.argsort(-sims[row_indices, idxs_part], axis=1)]
        idxs = idxs_sorted
        dists = -np.take_along_axis(sims, idxs, axis=1).astype(np.float32)

        if k_eff < k:
            pad_count = k - k_eff
            dists = np.pad(dists, ((0, 0), (0, pad_count)), constant_values=np.finfo(np.float32).max)
            idxs = np.pad(idxs, ((0, 0), (0, pad_count)), constant_values=-1)

        return dists, idxs.astype(np.int64)


class VectorIndex:
    def __init__(self, embeddings: Optional[np.ndarray], texts: Optional[List[str]] = None) -> None:
        texts = texts or []
        self.texts = list(texts)
        if embeddings is None or getattr(embeddings, "size", 0) == 0:
            self.index = VectorIndexFallback(np.zeros((0, 0), dtype=np.float32), self.texts)
            self._use_faiss = False
            return
        emb = np.asarray(embeddings, dtype=np.float32, order="C")
        if emb.ndim == 1:
            emb = emb.reshape(1, -1)
        self.embeddings = emb
        n, d = emb.shape
        if _HAS_FAISS and n > 0 and faiss is not None:
            try:
                normed = _safe_normalize(self.embeddings, axis=1)
                normed = np.ascontiguousarray(normed.astype(np.float32))
                self._faiss_index = faiss.IndexFlatIP(d)
                self._faiss_index.add(normed)
                self._use_faiss = True
            except Exception as ex:
                logger.exception("Failed to build FAISS index; falling back to numpy. Error: %s", ex)
                self.index = VectorIndexFallback(self.embeddings, texts)
                self._use_faiss = False
        else:
            self.index = VectorIndexFallback(self.embeddings, texts)
            self._use_faiss = False

    def search(self, q_emb: np.ndarray, k: int = 5) -> Tuple[np.ndarray, np.ndarray]:
        q = np.asarray(q_emb, dtype=np.float32, order="C")
        q = _ensure_2d(q)
        n_q, dim_q = q.shape
        if getattr(self, "_use_faiss", False):
            if dim_q != self._faiss_index.d:
                logger.error("Query dim (%d) != FAISS index dim (%d)", dim_q, self._faiss_index.d)
                return np.full((n_q, k), np.finfo(np.float32).max, dtype=np.float32), np.full((n_q, k), -1, dtype=np.int64)
            q_normed = _safe_normalize(q, axis=1)
            try:
                D, I = self._faiss_index.search(np.ascontiguousarray(q_normed), k)
                D = D.astype(np.float32)
                dists = -D
                idxs = I.astype(np.int64)
                return dists, idxs
            except Exception as ex:
                logger.exception("FAISS search failed; falling back to numpy. Error: %s", ex)
                return self.index.search(q, k)
        else:
            return self.index.search(q, k)

# ------------------------------
# Embedding helpers & caching
# ------------------------------

if st is not None:
    def _cache_decorator(func):
        return st.cache_resource(func)
else:
    from functools import lru_cache
    def _cache_decorator(func):
        return lru_cache(maxsize=2)(func)


@_cache_decorator
def load_sentence_transformer(model_name: str = "all-MiniLM-L6-v2"):
    if not _HAS_SENTENCE_TRANSFORMERS or SentenceTransformer is None:
        raise RuntimeError("sentence-transformers is required for embeddings. Install with: pip install sentence-transformers")
    logger.info("Loading SentenceTransformer model '%s' ...", model_name)
    model = SentenceTransformer(model_name)
    return model


def embed_texts(model: "SentenceTransformer", texts: List[str]) -> np.ndarray:
    texts = texts or []
    if not texts:
        try:
            dim = model.get_sentence_embedding_dimension()
            return np.zeros((0, int(dim)), dtype=np.float32)
        except Exception:
            return np.zeros((0, 0), dtype=np.float32)
    arr = model.encode(texts, show_progress_bar=False, convert_to_numpy=True)
    arr = np.asarray(arr, dtype=np.float32)
    if arr.ndim == 1:
        arr = arr.reshape(1, -1)
    return arr

# ------------------------------
# Parsers: pptx & pdf
# ------------------------------

def extract_from_pptx_bytes(file_bytes: bytes) -> List[Dict]:
    """
    Extract slide-level text and embedded images from a PPTX file.
    Returns list of dicts: {"index": int, "text": str, "images": List[bytes]}
    Robust against unusual shapes and missing python-pptx.
    """
    if not _HAS_PPTX or Presentation is None:
        return [{"index": 0, "text": "[python-pptx not installed]", "images": []}]

    slides: List[Dict] = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for i, slide in enumerate(prs.slides):
            texts: List[str] = []
            images: List[bytes] = []
            for shape in slide.shapes:
                # Text extraction (robust)
                try:
                    if getattr(shape, "has_text_frame", False):
                        tf = getattr(shape, "text_frame", None)
                        if tf:
                            # concatenate non-empty paragraphs
                            pts = [p.text.strip() for p in tf.paragraphs if getattr(p, "text", None) and p.text.strip()]
                            if pts:
                                texts.append("\n".join(pts))
                except Exception:
                    logger.debug("Skipping pptx shape text extraction on slide %d", i, exc_info=True)

                # Image extraction (handles picture shapes)
                try:
                    # python-pptx picture shapes often have .image.blob
                    img = getattr(shape, "image", None)
                    if img is not None:
                        blob = getattr(img, "blob", None)
                        if blob:
                            images.append(blob)
                except Exception:
                    # some shapes may provide image via shape._pic or anchor; skip gracefully
                    logger.debug("Skipping pptx shape image extraction on slide %d", i, exc_info=True)

            slides.append({"index": i, "text": "\n".join(texts).strip(), "images": images})
        return slides
    except Exception as e:
        logger.exception("pptx parse error: %s", e)
        return [{"index": 0, "text": f"[pptx parse error] {e}", "images": []}]



def extract_from_pdf_bytes(file_bytes: bytes) -> List[Dict]:
    """
    Extract per-page text and images from a PDF (PyMuPDF). If a page has no selectable text
    and no embedded images, render the page to an image (PNG bytes) and return that as the page image.
    """
    if not _HAS_PYMUPDF or fitz is None:
        return [{"index": 0, "text": "[pymupdf not installed, can't parse PDF]", "images": []}]

    slides: List[Dict] = []
    doc = None
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        page_count = getattr(doc, "page_count", None) or doc.pageCount if hasattr(doc, "pageCount") else len(doc)
        for i in range(page_count):
            page = doc.load_page(i)
            # extract selectable text if any
            try:
                text = page.get_text("text") or ""
                text = text.strip()
            except Exception:
                logger.debug("Page %d: text extraction failed; treating as empty.", i, exc_info=True)
                text = ""

            images: List[bytes] = []
            # try to extract embedded images
            try:
                for img in page.get_images(full=True):
                    xref = img[0]
                    try:
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image.get("image")
                        if image_bytes:
                            images.append(image_bytes)
                    except Exception:
                        logger.debug("Failed to extract image xref %s on page %d", xref, i, exc_info=True)
            except Exception:
                logger.debug("get_images failed on page %d", i, exc_info=True)

            # If page has neither selectable text nor embedded images, rasterize page to PNG
            if (not text) and (not images):
                try:
                    mat = fitz.Matrix(2.0, 2.0)  # raster scale (2x)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    # PyMuPDF versions differ on pix.tobytes signature
                    try:
                        img_bytes = pix.tobytes(output="png")
                    except TypeError:
                        img_bytes = pix.tobytes()
                    images.append(img_bytes)
                except Exception as e:
                    logger.debug("PDF page rendering failed for page %d: %s", i, e, exc_info=True)

            slides.append({"index": i, "text": text, "images": images})
        return slides
    except Exception as e:
        logger.exception("pdf parse error: %s", e)
        return [{"index": 0, "text": f"[pdf parse error] {e}", "images": []}]
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass


# ------------------------------
# EasyOCR reader & OCR
# ------------------------------

@_cache_decorator
def get_easyocr_reader(lang_list: List[str] = ["en"]):
    """
    Return an EasyOCR reader instance or None if easyocr not available.
    Cached to avoid re-initializing repeatedly.
    """
    if not _HAS_EASYOCR or easyocr is None:
        logger.info("EasyOCR not available in this environment.")
        return None

    try:
        # allow passing list or comma-separated string
        langs = lang_list if isinstance(lang_list, (list, tuple)) else [l.strip() for l in str(lang_list).split(",") if l.strip()]
        reader = easyocr.Reader(langs, gpu=False)
        return reader
    except Exception as e:
        logger.warning("easyocr init failed: %s", e)
        return None



def ocr_image_bytes_list(image_bytes_list: List[bytes]) -> List[str]:
    """
    Run OCR on a list of image bytes and return list of extracted texts (one per image).
    Returns empty string for pages where OCR fails or reader not available.
    """
    results: List[str] = []
    reader = get_easyocr_reader(["en"])
    for idx, b in enumerate(image_bytes_list or []):
        if not b:
            results.append("")
            continue
        try:
            img = Image.open(io.BytesIO(b)).convert("RGB")
        except Exception as e:
            logger.debug("ocr_image_bytes_list: PIL open failed for image %d: %s", idx, e)
            results.append("")
            continue

        if _HAS_EASYOCR and reader is not None:
            try:
                arr = np.array(img)
                # easyocr Reader.readtext signature changed historically: try paragraph=True then fallback
                try:
                    raw = reader.readtext(arr, detail=0, paragraph=True)
                except TypeError:
                    raw = reader.readtext(arr, detail=0)
                if isinstance(raw, list):
                    txt = "\n".join([r for r in raw if r])
                else:
                    txt = str(raw)
                results.append(txt.strip())
            except Exception as e:
                logger.debug("easyocr readtext failed for image %d: %s", idx, e, exc_info=True)
                results.append("")
        else:
            # no OCR available
            results.append("")
    return results


# ------------------------------
# Text chunker (user-facing helper)
# ------------------------------

def chunk_text(text: str, max_chars: int = 700) -> List[str]:
    if not text:
        return []
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    parts = []
    cur = ""
    for p in paragraphs:
        if not cur:
            cur = p
        elif len(cur) + 1 + len(p) <= max_chars:
            cur = cur + "\n" + p
        else:
            parts.append(cur)
            cur = p
    if cur:
        parts.append(cur)
    final_parts = []
    for part in parts:
        if len(part) <= max_chars:
            final_parts.append(part)
        else:
            sents = part.split(". ")
            cur2 = ""
            for s in sents:
                piece = s if s.endswith('.') else s + '.'
                if not cur2:
                    cur2 = piece
                elif len(cur2) + 1 + len(piece) <= max_chars:
                    cur2 = cur2 + ' ' + piece
                else:
                    final_parts.append(cur2.strip())
                    cur2 = piece
            if cur2:
                final_parts.append(cur2.strip())
    return final_parts

# ------------------------------
# OpenRouter / Chat helpers
# ------------------------------

def call_openrouter_chat(system_prompt: str, user_prompt: str, model: str = "gpt-4o-mini", max_tokens: int = 512, temperature: float = 0.2) -> str:
    api_key = st.session_state.get("OPENROUTER_API_KEY") or os.getenv("OPENROUTER_API_KEY", DEFAULT_OPENROUTER_KEY)
    if not api_key:
        return "[OpenRouter error] No API key configured. Set it in Settings or env var OPENROUTER_API_KEY."
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    body = {"model": model, "messages": [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}], "max_tokens": max_tokens, "temperature": temperature}
    text = ""
    try:
        resp = requests.post(OPENROUTER_API_URL, headers=headers, json=body, timeout=30)
        text = getattr(resp, "text", "")
        resp.raise_for_status()
        data = resp.json()
        choices = data.get("choices") or []
        if choices:
            first = choices[0]
            msg = first.get("message") or {}
            content = msg.get("content") if isinstance(msg, dict) else None
            if not content:
                content = first.get("text")
            if not content:
                content = json.dumps(first)
            return content
        if "result" in data:
            return str(data["result"])
        return "[openrouter] empty response"
    except Exception as e:
        try:
            return f"[OpenRouter error] {e} — resp_text={text}"
        except Exception:
            return f"[OpenRouter error] {e}"


def extract_json_from_text(text: str) -> Optional[Any]:
    if not text:
        return None
    cleaned = text.strip()
    for fence in ['```json', '```']:
        cleaned = cleaned.replace(fence, '')
    import re
    m = re.search(r'([\{\[])', cleaned)
    if not m:
        return None
    start = m.start(1)
    for end in range(len(cleaned), start, -1):
        try:
            candidate = cleaned[start:end]
            parsed = json.loads(candidate)
            return parsed
        except Exception:
            continue
    arr_match = re.search(r'\[.*\]', cleaned, flags=re.S)
    if arr_match:
        try:
            return json.loads(arr_match.group(0))
        except Exception:
            pass
    obj_match = re.search(r'\{.*\}', cleaned, flags=re.S)
    if obj_match:
        try:
            return json.loads(obj_match.group(0))
        except Exception:
            pass
    return None

# ------------------------------
# Generators & utilities
# ------------------------------

def generate_multilevel_lesson(slide_text: str, related_texts: str = "") -> str:
    system = (
        "You are a patient expert teacher. Produce three sections titled 'Beginner', 'Intermediate', and 'Advanced'. "
        "For each section: give a clear explanation, one small worked example, and 2-3 concise tips. Finally, produce "
        "3 multiple-choice questions (with answers) labelled under 'MCQs:'. Keep the tone friendly and short."
    )
    user_prompt = f"SLIDE_TEXT:\n{slide_text}\n\nRELATED:\n{related_texts}\n\nRespond with sections as requested."
    return call_openrouter_chat(system, user_prompt, max_tokens=900, temperature=0.15)


def generate_deep_lesson(slide_text: str, related_texts: str = "") -> str:
    system = (
        "You are an expert university instructor teaching the topic in depth to a BTech 2nd-year mechanical engineering student. "
        "Start from first principles in very simple language, then progressively increase difficulty. For each level include: (1) concept explanation, "
        "(2) 2 worked examples (one numerical, one conceptual), (3) common mistakes, (4) 5 short practice MCQs with answers, and (5) 8 concise flashcard-style Q/A pairs. "
        "Label sections clearly: 'Foundations', 'Developing', 'Advanced'. Keep numerical steps explicit and provide units where applicable."
    )
    user_prompt = f"SLIDE_TEXT:\n{slide_text}\n\nRELATED:\n{related_texts}\n\nRespond with detailed lesson as requested."
    return call_openrouter_chat(system, user_prompt, model="gpt-4o-mini", max_tokens=2500, temperature=0.05)


def generate_mcq_set_from_text(text: str, qcount: int = 5) -> List[Dict]:
    system = "You are an assistant that generates high-quality multiple-choice questions from source text. Reply JSON only."
    user_prompt = (f"Create {qcount} multiple-choice questions based solely on the following text. For each question provide: "
                   f"'question' (string), 'options' (array of 4 strings), and 'answer_index' (0-based integer). Reply ONLY with a JSON array.\n\nTEXT:\n{text}\n")
    resp = call_openrouter_chat(system, user_prompt, max_tokens=700, temperature=0.0)
    parsed = extract_json_from_text(resp)
    if isinstance(parsed, list):
        return parsed
    return [{"question": "Summarize the main idea in one sentence.", "options": ["A", "B", "C", "D"], "answer_index": 0}]


def generate_flashcards_from_text(text: str, n: int = 10) -> List[Dict]:
    system = "Extract concise Q/A flashcards from the source text. Reply only in JSON array."
    user_prompt = f"Extract up to {n} short Q/A flashcards from the text below. Keep Q under 80 chars and A under 200 chars. Reply ONLY in JSON.\n\nTEXT:\n{text}"
    resp = call_openrouter_chat(system, user_prompt, max_tokens=700, temperature=0.0)
    parsed = extract_json_from_text(resp)
    if isinstance(parsed, list):
        return parsed
    return []

# ------------------------------
# SM-2 spaced repetition update
# ------------------------------

def sm2_update_card(easiness: float, interval: int, repetitions: int, quality: int) -> Tuple[float, int, int, int]:
    """
    SM-2 algorithm update. Returns (easiness, interval, repetitions, next_review_timestamp).
    """
    try:
        q = int(quality)
    except Exception:
        q = 0

    eas = float(easiness or 2.5)
    inter = int(interval or 1)
    reps = int(repetitions or 0)

    if q < 3:
        reps = 0
        inter = 1
    else:
        reps += 1
        if reps == 1:
            inter = 1
        elif reps == 2:
            inter = 6
        else:
            inter = max(1, int(round(inter * eas)))

    eas = max(1.3, eas + 0.1 - (5 - q) * (0.08 + (5 - q) * 0.02))
    next_review = int(time.time()) + int(inter) * 24 * 3600
    return eas, inter, reps, next_review


# ------------------------------
# Export & TTS helpers
# ------------------------------

def anki_export_csv_for_upload(upload_id: int, conn: sqlite3.Connection) -> Optional[Tuple[str, bytes]]:
    """
    Export flashcards for a given upload_id as TSV suitable for Anki import.
    Returns (filename, bytes) or None on error.
    """
    if conn is None:
        logger.debug("anki_export_csv_for_upload: no DB connection provided.")
        return None
    try:
        c = conn.cursor()
        c.execute("SELECT question, answer FROM flashcards WHERE upload_id = ?", (upload_id,))
        rows = c.fetchall() or []
        lines = []
        for q, a in rows:
            q2 = (q or "").replace("\t", " ").replace("\n", " ")
            a2 = (a or "").replace("\t", " ").replace("\n", " ")
            lines.append(f"{q2}\t{a2}")
        csv_bytes = "\n".join(lines).encode("utf-8")
        fname = f"slide_tutor_upload_{upload_id}_flashcards.txt"
        return fname, csv_bytes
    except Exception as e:
        logger.exception("anki_export_csv_for_upload failed: %s", e)
        return None



def text_to_speech_download(text: str, lang: str = "en") -> Tuple[str, bytes]:
    """
    Use gTTS to synthesize text to MP3 bytes. Raises RuntimeError with helpful message if not available.
    """
    if not _HAS_GTTS or gTTS is None:
        raise RuntimeError("gTTS not installed (pip install gTTS) or not available in this environment.")
    if not text:
        raise ValueError("No text provided for TTS.")
    try:
        tts = gTTS(text=text, lang=lang)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        tmp.close()
        tts.save(tmp.name)
        with open(tmp.name, "rb") as fh:
            b = fh.read()
        try:
            os.unlink(tmp.name)
        except Exception:
            pass
        return "lesson_audio.mp3", b
    except Exception as e:
        logger.exception("text_to_speech_download failed: %s", e)
        raise


# ------------------------------
# Streamlit UI
# ------------------------------

st.set_page_config(page_title=globals().get("APP_TITLE", "SlideTutor"), layout="wide")

# session state: API key
if "OPENROUTER_API_KEY" not in st.session_state:
    st.session_state["OPENROUTER_API_KEY"] = os.getenv("OPENROUTER_API_KEY", DEFAULT_OPENROUTER_KEY)

# CSS + wrapper
st.markdown("""
    <style>
    :root{ --bg:#071025; --bg-2:#06101A; --text:#E6F0FA; --muted:#9AA6B2; --accent:#2AB7A9; --accent-2:#4D7CFE; --radius:12px; --gap-md:20px; }
    .app-theme, .app-theme * { box-sizing: border-box; font-family: "Inter", system-ui, -apple-system, "Segoe UI", Roboto, "Helvetica Neue", Arial; color:var(--text); }
    .app-theme .container{ padding:20px; min-height:100vh; }
    .app-theme .topbar{ display:flex; gap:12px; align-items:center; justify-content:space-between; padding:12px 18px; }
    .app-theme .logo{ width:48px;height:48px;border-radius:12px; display:flex;align-items:center;justify-content:center;font-weight:700;color:white;background:linear-gradient(135deg,var(--accent-2),var(--accent)); box-shadow:0 12px 30px rgba(0,0,0,0.45); }
    .app-theme .title{ font-size:1.2rem; font-weight:700; margin:0 }
    .app-theme .subtitle{ color:var(--muted); font-size:0.88rem; margin-top:2px }
    .app-theme .tabs{ display:flex; gap:8px; align-items:center; }
    .app-theme .tab{ padding:8px 12px; border-radius:10px; font-weight:600; cursor:pointer; border:1px solid transparent; color:var(--muted); background:transparent; }
    .app-theme .tab.active{ background: linear-gradient(180deg, rgba(255,255,255,0.02), rgba(255,255,255,0.01)); color:var(--text); box-shadow:0 10px 28px rgba(0,0,0,0.45); }
    .app-theme .card{ background: linear-gradient(180deg, rgba(255,255,255,0.02), rgba(255,255,255,0.01)); padding:18px; border-radius:12px; border:1px solid rgba(255,255,255,0.03); box-shadow: 0 8px 28px rgba(2,6,12,0.5); }
    .app-theme .small-muted{ color:var(--muted); font-size:13px; }
    .app-theme .blob{ position: fixed; pointer-events:none; filter: blur(36px); z-index:0; opacity:0.55; }
    .app-theme .blob.one{ width:420px;height:420px;border-radius:50%; left:-140px; top:-120px; background: radial-gradient(circle at 30% 30%, rgba(77,124,254,0.18), transparent 30%); }
    .app-theme .blob.two{ width:320px;height:320px;border-radius:50%; right:-80px; bottom:-80px; background: radial-gradient(circle at 70% 70%, rgba(42,183,169,0.12), transparent 30%); }
    .app-theme .small-actions{ display:flex; gap:8px; }
    .app-theme .slide-img{ max-width:100%; border-radius:8px; box-shadow: 0 8px 20px rgba(0,0,0,0.45); }
    .app-theme .muted{ color:var(--muted) }
    </style>
    <div class="app-theme">
""", unsafe_allow_html=True)

st.markdown("<div class='blob one'></div><div class='blob two'></div>", unsafe_allow_html=True)

with st.container():
    top_cols = st.columns([0.4, 3, 1.2])
    with top_cols[0]:
        st.markdown("<div class='logo'>ST</div>", unsafe_allow_html=True)
    with top_cols[1]:
        app_title = globals().get("APP_TITLE", "SlideTutor")
        app_sub = globals().get("APP_SUBTITLE", "AI slides → lessons • quizzes • flashcards")
        st.markdown(f"<div class='title'>{app_title}</div><div class='subtitle'>{app_sub}</div>", unsafe_allow_html=True)
    with top_cols[2]:
        st.markdown("<div style='text-align:right'><small class='small-muted'>Student edition • Improved</small></div>", unsafe_allow_html=True)

TAB_NAMES = ["Home", "Upload & Process", "Lessons", "Chat Q&A", "Quizzes", "Flashcards", "Export", "Progress", "Settings"]
if "active_tab" not in st.session_state:
    st.session_state["active_tab"] = "Upload & Process"

# Render top tabs
tab_cols = st.columns(len(TAB_NAMES))
for i, name in enumerate(TAB_NAMES):
    pressed = tab_cols[i].button(name, key=f"tab_btn_{name}")
    if pressed:
        st.session_state["active_tab"] = name

active_tab = st.session_state["active_tab"]

if "uploads" not in st.session_state:
    st.session_state["uploads"] = []

# helpers

def _simple_chunk_text(text: str, chunk_size: int = 800, overlap: int = 80) -> List[str]:
    """
    Split text into roughly chunk_size-character segments with 'overlap' characters overlap.
    Keeps paragraphs together when possible.
    """
    if not text:
        return []
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks: List[str] = []
    for p in paragraphs:
        if len(p) <= chunk_size:
            chunks.append(p)
            continue
        start = 0
        L = len(p)
        while start < L:
            end = min(L, start + chunk_size)
            chunk = p[start:end].strip()
            chunks.append(chunk)
            if end >= L:
                break
            start = max(0, end - overlap)
    return chunks


def _make_upload_id(filename: str) -> str:
    """Return a short pseudo-unique upload id based on filename + timestamp hash."""
    import hashlib
    h = hashlib.sha1((filename + str(time.time())).encode("utf-8")).hexdigest()[:12]
    safe_name = filename.replace(" ", "_")
    return f"{h}_{safe_name}"


def add_upload_file(uploaded_file) -> Dict:
    """
    Process UploadedFile-like into upload dict:
    - extracts slides (pptx/pdf)
    - runs per-slide chunking (uses chunk_text() if present)
    - builds mapping list and 'full_text'
    """
    try:
        fb = uploaded_file.read()
    except Exception as e:
        st.error(f"Failed to read uploaded file: {e}")
        return {}

    filename = getattr(uploaded_file, "name", None) or getattr(uploaded_file, "filename", None) or "<uploaded_file>"
    ext = filename.lower().split(".")[-1]
    kind = "pdf" if ext == "pdf" else "pptx" if ext in ("pptx", "ppt") else "unknown"

    # extract slides
    slides: List[Dict] = []
    try:
        if kind == "pptx":
            slides = extract_from_pptx_bytes(fb)
        elif kind == "pdf":
            slides = extract_from_pdf_bytes(fb)
        else:
            # try pdf then pptx as fallback
            try:
                slides = extract_from_pdf_bytes(fb)
            except Exception:
                try:
                    slides = extract_from_pptx_bytes(fb)
                except Exception:
                    slides = [{"index": 0, "text": "[unsupported format]", "images": []}]
    except Exception as e:
        logger.exception("add_upload_file extraction failed: %s", e)
        slides = [{"index": 0, "text": f"[extraction failed] {e}", "images": []}]

    # build chunks and mapping (prefer user's chunk_text if available)
    chunks: List[str] = []
    mapping: List[Dict[str, Any]] = []
    for s in slides:
        sidx = int(s.get("index", 0))
        text = s.get("text", "") or ""
        try:
            if "chunk_text" in globals() and callable(globals()["chunk_text"]):
                parts = globals()["chunk_text"](text)
            else:
                parts = _simple_chunk_text(text)
        except Exception as e:
            logger.debug("chunking failed for slide %s: %s", sidx, e)
            parts = _simple_chunk_text(text)

        for p in parts:
            chunks.append(p)
            mapping.append({"slide": sidx, "text": p})

    # full_text: concatenated slide texts (useful for whole-doc operations)
    full_text = "\n\n".join([(s.get("text") or "").strip() for s in slides]).strip()
    upload = {
        "id": _make_upload_id(filename),
        "filename": filename,
        "bytes": fb,
        "kind": kind,
        "slides": slides,
        "slide_count": len(slides),
        "full_text": full_text,
        "chunks": chunks,
        "mapping": mapping,
        "embeddings": None,
        "index": None,
        "index_built": False,
        "status_msg": "Uploaded",
    }

    if "uploads" not in st.session_state:
        st.session_state["uploads"] = []
    st.session_state["uploads"].append(upload)
    return upload



def build_index_for_upload(upload: Dict, model_name: Optional[str] = None, use_full_text: bool = False) -> Dict:
    """
    Build embeddings and a VectorIndex for an upload.
    - If use_full_text is True, embed the concatenated full_text as a single vector.
    - Otherwise embed upload['chunks'] as usual.
    Returns the modified upload dict (in-place).
    """
    if upload is None:
        raise ValueError("upload is None")

    model_name = model_name or globals().get("EMBEDDING_MODEL_NAME", "all-MiniLM-L6-v2")

    # attempt to load model
    try:
        model = load_sentence_transformer(model_name)
    except Exception as e:
        upload["status_msg"] = f"Embedding model load failed: {e}"
        logger.debug(upload["status_msg"])
        st.error(upload["status_msg"])
        return upload

    # choose texts to embed
    if use_full_text:
        texts = []
        ft = upload.get("full_text", "") or ""
        if not ft:
            # fallback: join chunks
            texts = ["\n\n".join(upload.get("chunks", []) or [])]
        else:
            texts = [ft]
    else:
        texts = upload.get("chunks", []) or []

    if not texts:
        # create zero-sized embeddings with known dim if possible
        try:
            dim = model.get_sentence_embedding_dimension()
        except Exception:
            dim = 0
        if dim > 0:
            upload["embeddings"] = np.zeros((0, dim), dtype=np.float32)
        else:
            upload["embeddings"] = np.zeros((0, 0), dtype=np.float32)
        upload["index"] = VectorIndex(upload["embeddings"], [])
        upload["index_built"] = True
        upload["status_msg"] = "No texts to embed (empty)."
        st.success(upload["status_msg"])
        return upload

    # compute embeddings
    try:
        with st.spinner("Computing embeddings..."):
            arr = embed_texts(model, texts)
            if arr is None:
                raise RuntimeError("embed_texts returned None")
            arr = np.asarray(arr, dtype=np.float32)
            upload["embeddings"] = arr
            # If use_full_text we still want index to reference texts (single entry)
            upload["index"] = VectorIndex(arr, texts)
            upload["index_built"] = True
            upload["status_msg"] = f"Index built ({arr.shape[0]} vectors, dim={arr.shape[1] if arr.ndim > 1 else 'unknown'})"
            st.success(upload["status_msg"])
            return upload
    except Exception as ex:
        upload["status_msg"] = f"Index build failed: {ex}"
        logger.exception("build_index_for_upload failed: %s", ex)
        st.error(upload["status_msg"])
        # fallback to empty index so other code doesn't crash
        upload["embeddings"] = np.zeros((0, 0), dtype=np.float32)
        upload["index"] = VectorIndex(upload["embeddings"], upload.get("chunks", []))
        upload["index_built"] = False
        return upload



def download_upload_json(upload: Dict) -> None:
    """
    Render a data: URL download link for the upload metadata (slides & chunks).
    """
    if not upload:
        st.warning("No upload provided for download.")
        return
    safe = {
        "filename": upload.get("filename"),
        "kind": upload.get("kind"),
        "slide_count": upload.get("slide_count", len(upload.get("slides", []))),
        "slides": [{"index": s.get("index"), "text": s.get("text"), "images": len(s.get("images", []))} for s in upload.get("slides", [])],
        "chunks_count": len(upload.get("chunks", [])),
    }
    payload = json.dumps(safe, indent=2, ensure_ascii=False)
    b64 = base64.b64encode(payload.encode()).decode()
    href = f'<a href="data:application/json;base64,{b64}" download="{upload.get("filename")}_meta.json">Download metadata (JSON)</a>'
    st.markdown(href, unsafe_allow_html=True)


# ---------- Home tab ----------
if active_tab := st.session_state.get("active_tab") == "Home":
    st.markdown("<div class='card' style='padding:18px'>", unsafe_allow_html=True)
    c1, c2 = st.columns([3, 1])
    with c1:
        st.markdown(f"<h2 style='margin:0'>Welcome to <span style='color:#a79bff'>{APP_TITLE}</span></h2>", unsafe_allow_html=True)
        st.markdown(f"<div class='small-muted'>Quickly upload slide decks (PPTX/PDF), build searchable semantic indices, and auto-generate lessons, quizzes and flashcards.</div>", unsafe_allow_html=True)
        
        st.markdown("<div class='small-muted'>Highlights:</div>", unsafe_allow_html=True)
        st.markdown("<ul class='small-muted'><li>Scanned PDFs supported (raster fallback)</li><li>On-demand embedding / search (FAISS if available)</li><li>Export metadata & previews</li></ul>", unsafe_allow_html=True)
        st.markdown("<div class='small-actions'>", unsafe_allow_html=True)
        if st.button("Upload & Process"):
            st.session_state["active_tab"] = "Upload & Process"
        if st.button("Practice Flashcards"):
            st.session_state["active_tab"] = "Flashcards"
        if st.button("Generate Lesson"):
            st.session_state["active_tab"] = "Lessons"
        st.markdown("</div>", unsafe_allow_html=True)
    with c2:
        st.markdown("<div class='card' style='padding:12px;text-align:center'>", unsafe_allow_html=True)
        st.markdown("<h4 style='margin:6px 0'>Usage tip</h4>", unsafe_allow_html=True)
        st.markdown("<div class='small-muted'>Start by uploading your .pptx or .pdf in the 'Upload & Process' tab. Build the embedding index to enable Search, Lessons and Quizzes.</div>", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# ---------- Upload & Process tab ----------
if st.session_state.get("active_tab") == "Upload & Process":
    st.markdown("<div class='card' style='margin-bottom:12px'>", unsafe_allow_html=True)
    st.header("Upload & Process")
    st.markdown("<div class='small-muted'>Upload one or more PPTX / PDF files. Each upload is processed into slide-level text + images and chunked into semantically meaningful pieces.</div>", unsafe_allow_html=True)

    uploaded = st.file_uploader("Choose PPTX/PDF files", accept_multiple_files=True, type=["pdf", "pptx", "ppt"])
    if uploaded:
        for f in uploaded:
            known_names = {u.get("filename") for u in st.session_state.get("uploads", [])}
            fname = getattr(f, "name", None) or "<uploaded_file>"
            if fname in known_names:
                st.warning(f"File '{fname}' already uploaded - skipping duplicate.")
                continue
            with st.spinner(f"Processing {fname} ..."):
                up = add_upload_file(f)
                if up and isinstance(up, dict):
                    uploaded_name = up.get("filename") or fname
                    st.success(f"Uploaded: {uploaded_name}")
                elif up:
                    st.success(f"Uploaded: {fname}")
                else:
                    st.error(f"Failed to process upload: {fname}")
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div style='margin-top:12px'>", unsafe_allow_html=True)
    uploads = st.session_state.get("uploads", [])
    if not uploads:
        st.info("No uploads yet. Use the uploader above to add a PPTX or PDF.")
    else:
        for idx, up in enumerate(list(uploads)):
            with st.expander(f"{up['filename']} — {up.get('status_msg','')}"):
                cols = st.columns([3, 1])
                with cols[0]:
                    st.markdown(f"**File:** {up['filename']} • **Type:** {up['kind']}  ")
                    st.markdown(f"<div class='small-muted'>Chunks: {len(up.get('chunks',[]))} • Index built: {up.get('index_built')}</div>", unsafe_allow_html=True)
                    slides = up.get("slides", [])[:4]
                    for s in slides:
                        st.markdown(f"**Slide {s.get('index', '?')}**")
                        txt = s.get("text", "").strip()
                        if txt:
                            st.markdown(f"<div class='small-muted'>{txt[:800]}{'...' if len(txt)>800 else ''}</div>", unsafe_allow_html=True)
                        images = s.get("images", []) or []
                        if images:
                            try:
                                st.image(images[0], caption=f"Slide {s.get('index')} preview", use_container_width=True)
                            except Exception:
                                st.markdown("<div class='small-muted'>[image preview not renderable]</div>", unsafe_allow_html=True)
                        st.markdown("---")
                with cols[1]:
                    if st.button("Build / Rebuild Index", key=f"build_idx_{up['id']}"):
                        st.session_state["uploads"][idx] = build_index_for_upload(up)
                    if st.button("Download metadata", key=f"dlmeta_{up['id']}"):
                        download_upload_json(up)
                    if st.button("Delete", key=f"del_{up['id']}"):
                        st.session_state["uploads"].pop(idx)
                        st.experimental_rerun()
                    
                    st.markdown(f"<div class='small-muted'>{up.get('status_msg','')}</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# ---------- Lessons tab ----------
if st.session_state.get("active_tab") == "Lessons":
    st.header("Lessons (Generate from uploaded files)")
    uploads = st.session_state.get("uploads", [])
    if not uploads:
        st.info("Upload a file first (Upload & Process tab).")
    else:
        options = [f"{u['filename']} ({'indexed' if u.get('index_built') else 'not indexed'})" for u in uploads]
        sel = st.selectbox("Select upload to generate lessons from", options=options)
        sel_idx = options.index(sel)
        selected = uploads[sel_idx]
        st.markdown("<div class='small-muted'>Choose a generation mode and press Generate.</div>", unsafe_allow_html=True)
        mode = st.radio("Generation Mode", ["Concise lesson", "Detailed lesson with examples", "Key-points summary"], index=0)
        if st.button("Generate Lesson"):
            if not selected.get("index_built"):
                st.error("Index not built for this upload. Please go to 'Upload & Process' and press 'Build / Rebuild Index'.")
            elif not st.session_state.get("OPENROUTER_API_KEY"):
                st.error("No OPENROUTER_API_KEY set. Set it in Settings (or as environment variable).")
            else:
                if "generate_lesson_from_upload" in globals():
                    try:
                        with st.spinner("Generating lesson..."):
                            lesson_text = globals()["generate_lesson_from_upload"](selected, mode=mode)
                            st.success("Lesson generated")
                            st.markdown(lesson_text)
                    except Exception as e:
                        st.error(f"Lesson generation failed: {e}")
                else:
                    st.warning("No lesson-generation pipeline found in this runtime. Implement `generate_lesson_from_upload(upload, mode)` to enable this feature.")

# ---------- Chat Q&A, Quizzes, Flashcards, Export, Progress, Settings ----------
if st.session_state.get("active_tab") == "Chat Q&A":
    st.header("Chat Q&A")
    st.markdown("Ask a question based on your uploaded slides (requires index & API key).")
    q = st.text_input("Enter question")
    if st.button("Ask"):
        uploads = [u for u in st.session_state.get("uploads", []) if u.get("index_built")]
        if not uploads:
            st.error("No indexed uploads available. Build an index first.")
        else:
            if "answer_question" in globals():
                try:
                    res = globals()["answer_question"](q, uploads)
                    st.markdown(res)
                except Exception as e:
                    st.error(f"Q&A failed: {e}")
            else:
                st.info("Q&A pipeline not available. Implement answer_question(query, indexed_uploads).")

if st.session_state.get("active_tab") == "Quizzes":
    st.header("Auto-generated Quizzes")
    st.markdown("Generate multiple-choice quizzes from selected upload (placeholder).")
    if st.button("Generate Quiz (placeholder)"):
        st.info("Implement quiz generation or call your existing function (e.g., generate_quiz_from_upload).")

if st.session_state.get("active_tab") == "Flashcards":
    st.header("Flashcards / Spaced Repetition")
    st.markdown("Practice flashcards generated from content. This is a placeholder UI; wire it to your SRS backend.")
    if st.button("Practice now"):
        st.info("Practice logic not wired. Implement or connect to your SRS (SM-2) implementation.")

if st.session_state.get("active_tab") == "Export":
    st.header("Export")
    st.markdown("Export all metadata and optionally embeddings (if index built).")
    if st.button("Export all metadata JSON"):
        exports = []
        for u in st.session_state.get("uploads", []):
            exports.append({"filename": u.get("filename"), "kind": u.get("kind"), "chunks_count": len(u.get("chunks", [])), "indexed": bool(u.get("index_built"))})
        payload = json.dumps(exports, indent=2)
        b64 = base64.b64encode(payload.encode()).decode()
        st.markdown(f'<a href="data:application/json;base64,{b64}" download="slide_tutor_export.json">Download export</a>', unsafe_allow_html=True)

if st.session_state.get("active_tab") == "Progress":
    st.header("Progress")
    st.markdown("Track your learning progress here (placeholder).")

if st.session_state.get("active_tab") == "Settings":
    st.header("Settings")
    st.text_input("OpenRouter API Key (session only)", value=st.session_state.get("OPENROUTER_API_KEY", ""), key="OPENROUTER_API_KEY_INPUT", type="password")
    if st.button("Save API key to session"):
        st.session_state["OPENROUTER_API_KEY"] = st.session_state.get("OPENROUTER_API_KEY_INPUT")
        st.success("API key saved to session.")
    st.markdown("<div class='small-muted'>Tip: Prefer storing API keys as environment variables in production.</div>", unsafe_allow_html=True)

# close wrapper
st.markdown("</div>", unsafe_allow_html=True)

# End of file

# --------------------------------------------------------------------------------
# The actual UI tabs (keeps exactly the features you had, but hardened)
# --------------------------------------------------------------------------------

# Upload & Process
# --------------------------
# Replacement UI: Upload & Process + Lessons / Chat / Quizzes / Flashcards / Export / Progress / Settings
# This block is a drop-in replacement starting at:
# if active_tab == "Upload & Process":
# and ending at the Settings tab handling.
#
# Assumes helper functions/objects exist elsewhere in the module:
# - extract_from_pdf_bytes, extract_from_pptx_bytes
# - _ocr_images_safe, _chunk_text_safe, _build_index_safe
# - VectorIndex, _get_db_cursor, _safe_call, upload_db_id, _db_conn
# - generate_multilevel_lesson/generate_deep_lesson/generate_mcq_set_from_text/generate_flashcards_from_text
# - sm2_update_card, anki_export_csv_for_upload, text_to_speech_download
# - logging via log(...) or logger
# --------------------------

import hashlib

# --------------------------
# Upload & Process tab
# --------------------------
# --------------------------
# Replacement UI block (Upload & Process ... Settings)
# Paste to replace from: `if active_tab == "Upload & Process":` down through the Settings tab.
# --------------------------
import hashlib

# Helper: session-safe incremental id when no DB id
def _make_session_id() -> str:
    return str(int(time.time() * 1000))

# --------------------------
# Upload & Process tab
# --------------------------
if active_tab == "Upload & Process":
    st.header("Upload PPTX / PDF (Student Upload)")

    uploaded_file = st.file_uploader(
        "Choose a PPTX or PDF file",
        type=["pptx", "pdf"],
        accept_multiple_files=False,
        key="uploader_file",
    )

    if uploaded_file is not None:
        # read bytes safely
        try:
            raw_bytes = uploaded_file.read()
        except Exception as e:
            st.error(f"Could not read uploaded file: {e}")
            raw_bytes = b""

        # basic metadata
        fname = getattr(uploaded_file, "name", "uploaded_file")
        try:
            file_size = int(getattr(uploaded_file, "size", len(raw_bytes)))
        except Exception:
            file_size = len(raw_bytes)

        # warn on very large uploads
        MB = 1024 * 1024
        if file_size > 100 * MB:
            st.warning("Large file (>100MB) detected. Processing may be slow or may run out of memory.")

        st.info(f"Processing {fname} ...")

        try:
            # dedupe using SHA1 of bytes (if available)
            file_hash = hashlib.sha1(raw_bytes).hexdigest() if raw_bytes else None

            # extract slides depending on file extension
            slides: List[Dict[str, Any]] = []
            lower = fname.lower()
            if lower.endswith(".pptx") or lower.endswith(".ppt"):
                if not globals().get("_HAS_PPTX", False):
                    st.error("python-pptx not installed; cannot parse PPTX. Install python-pptx to enable this feature.")
                    slides = [{"index": 0, "text": "[python-pptx not installed]", "images": []}]
                else:
                    try:
                        slides = extract_from_pptx_bytes(raw_bytes)
                    except Exception as e:
                        log("pptx extraction error:", e)
                        slides = [{"index": 0, "text": f"[pptx parse error] {e}", "images": []}]
            elif lower.endswith(".pdf"):
                if not globals().get("_HAS_PYMUPDF", False):
                    st.error("pymupdf not installed; cannot parse PDF. Install pymupdf to enable full PDF parsing.")
                    slides = [{"index": 0, "text": "[pymupdf not installed, can't parse PDF]", "images": []}]
                else:
                    try:
                        slides = extract_from_pdf_bytes(raw_bytes)
                    except Exception as e:
                        log("pdf extraction error:", e)
                        slides = [{"index": 0, "text": f"[pdf parse error] {e}", "images": []}]
            else:
                # best-effort fallback: try PDF then PPTX
                try:
                    slides = extract_from_pdf_bytes(raw_bytes)
                except Exception:
                    try:
                        slides = extract_from_pptx_bytes(raw_bytes)
                    except Exception:
                        slides = [{"index": 0, "text": "[Unsupported file type or extraction failed]", "images": []}]

            # normalize slides list
            normalized = []
            for i, s in enumerate(slides):
                normalized.append({
                    "index": int(s.get("index", i)),
                    "text": (s.get("text") or "").strip(),
                    "images": s.get("images") or [],
                })
            slides = normalized

            # OCR images into slide text (if helper available)
            with st.spinner("Running OCR on slide images (if any)..."):
                for si, s in enumerate(slides):
                    imgs = s.get("images") or []
                    if imgs:
                        try:
                            ocr_texts = _ocr_images_safe(imgs)
                            if ocr_texts:
                                appended = "\n".join([t for t in ocr_texts if t])
                                if appended:
                                    s["text"] = (s.get("text", "") + "\n\n" + appended).strip()
                        except Exception as e:
                            log(f"OCR failure on slide {si}:", e)

            # chunk text and build mapping
            chunks: List[str] = []
            mapping: List[Dict[str, Any]] = []
            for s in slides:
                try:
                    parts = _chunk_text_safe(s.get("text", "") or "")
                    if not parts:
                        parts = [s.get("text", "") or ""]
                except Exception as e:
                    log("chunk_text failed for slide", s.get("index"), e)
                    parts = [s.get("text", "") or ""]
                for p in parts:
                    chunks.append(p)
                    mapping.append({"slide": int(s.get("index", 0)), "text": p})

            # upload object id uses ms timestamp + short hash where available
            session_id = _make_session_id()
            uid = f"{session_id}_{(file_hash[:10] if file_hash else 'nohash')}"

            upload_obj: Dict[str, Any] = {
                "id": uid,
                "filename": fname,
                "uploaded_at": int(time.time()),
                "file_size": file_size,
                "file_hash": file_hash,
                "slides": slides,
                "chunks": chunks,
                "mapping": mapping,
                "embeddings": None,
                "index": None,
                "index_built": False,
                "status_msg": "Uploaded",
            }

            # build index (safe wrapper)
            with st.spinner("Creating embeddings and index (if sentence-transformers is available)..."):
                try:
                    upload_obj = _build_index_safe(upload_obj) or upload_obj
                    # normalize index_built flag
                    if upload_obj.get("index_built") is None:
                        upload_obj["index_built"] = bool(upload_obj.get("embeddings") is not None and getattr(upload_obj.get("embeddings"), "size", 0) > 0)
                    upload_obj["status_msg"] = upload_obj.get("status_msg", "Index build attempted")
                except Exception as e:
                    log("Index build failed:", e)
                    # fallback to an empty index object
                    try:
                        upload_obj["embeddings"] = np.zeros((0, 1), dtype=np.float32)
                        upload_obj["index"] = VectorIndex(upload_obj["embeddings"], upload_obj.get("chunks", []))
                    except Exception as ex_idx:
                        log("Fallback vector index creation failed:", ex_idx)
                        upload_obj["index"] = None
                    upload_obj["index_built"] = False
                    upload_obj["status_msg"] = "Index build failed (fallback index created)"

            # persist upload metadata in DB (best-effort)
            try:
                cur = _get_db_cursor()
                if cur:
                    meta = {"n_slides": len(slides), "n_chunks": len(chunks), "file_size": file_size, "file_hash": file_hash}
                    cur.execute("INSERT INTO uploads (filename, uploaded_at, meta) VALUES (?, ?, ?)",
                                (fname, upload_obj["uploaded_at"], json.dumps(meta)))
                    _db_conn.commit()
                    db_id = getattr(cur, "lastrowid", None)
                    upload_obj["db_id"] = int(db_id) if db_id is not None else None
                else:
                    upload_obj["db_id"] = None
            except Exception as e:
                log("Could not save upload metadata to DB:", e)
                upload_obj["db_id"] = None

            # add to session uploads with deduplication by content-hash
            try:
                st.session_state.setdefault("uploads", [])
                replaced = False
                if upload_obj.get("file_hash"):
                    for i, existing in enumerate(list(st.session_state["uploads"])):
                        if existing.get("file_hash") and existing.get("file_hash") == upload_obj.get("file_hash"):
                            st.session_state["uploads"][i] = upload_obj
                            replaced = True
                            break
                if not replaced:
                    st.session_state["uploads"].append(upload_obj)
            except Exception as e:
                log("Could not append upload to session:", e)
                try:
                    st.session_state.setdefault("uploads", []).append(upload_obj)
                except Exception as ex2:
                    log("Final append failed:", ex2)

            # ensure user stays on Upload & Process tab
            st.session_state["active_tab"] = "Upload & Process"

            st.success(f"Upload processed: {len(slides)} slides/pages, {len(chunks)} chunks.")

            # compact info + download
            info_cols = st.columns([3, 1])
            with info_cols[0]:
                st.markdown(f"**Filename:** {fname}  •  **Slides:** {len(slides)}  •  **Chunks:** {len(chunks)}")
            with info_cols[1]:
                try:
                    json_bytes = json.dumps(slides, ensure_ascii=False, indent=2).encode("utf-8")
                    st.download_button("Download extracted slides (JSON)", json_bytes, file_name=f"{fname}_extracted.json", mime="application/json", key=f"dljson_{upload_obj['id']}")
                except Exception as e:
                    log("Could not prepare download button:", e)

            # ---------- Preview mode and viewer ----------
            preview_key = f"view_mode_{upload_obj.get('db_id') or upload_obj.get('id')}"
            view_mode = st.radio("Preview mode", ["Text (chunks)", "Slide Viewer"], index=0, key=preview_key)

            if view_mode == "Text (chunks)":
                with st.expander("Preview first 25 chunks (expand to view)"):
                    for i, c in enumerate(chunks[:25]):
                        st.code((c[:1200] + ("..." if len(c) > 1200 else "")))
            else:
                # Viewer card
                st.markdown("<div class='card'>", unsafe_allow_html=True)
                uid = str(upload_obj.get("db_id") or upload_obj.get("id"))
                viewer_index_key = f"viewer_index_{uid}"
                slider_key = f"viewer_slider_{uid}"
                prev_key = f"viewer_prev_{uid}"
                next_key = f"viewer_next_{uid}"

                st.session_state.setdefault(viewer_index_key, 0)

                slides_list = upload_obj.get("slides", []) or []
                num_slides = len(slides_list)

                # Nav row
                nav_c1, nav_c2, nav_c3 = st.columns([1, 2, 1])
                with nav_c1:
                    if st.button("◀ Prev", key=prev_key):
                        st.session_state[viewer_index_key] = max(0, st.session_state[viewer_index_key] - 1)
                with nav_c2:
                    st.markdown(f"### Slide {st.session_state[viewer_index_key] + 1} / {max(1, num_slides)}")
                with nav_c3:
                    if st.button("Next ▶", key=next_key):
                        st.session_state[viewer_index_key] = min(max(0, num_slides - 1), st.session_state[viewer_index_key] + 1)

                # main area + sidebar
                main_col, side_col = st.columns([3, 1])
                cur_idx = int(st.session_state.get(viewer_index_key, 0))
                cur_idx = max(0, min(cur_idx, max(0, num_slides - 1)))
                st.session_state[viewer_index_key] = cur_idx

                cur_slide = slides_list[cur_idx] if num_slides > 0 else {"index": 0, "text": "", "images": []}
                imgs = cur_slide.get("images") or []

                # Main display: use_container_width replacement for deprecated parameter
                with main_col:
                    if imgs:
                        img_bytes = imgs[0]
                        try:
                            from PIL import Image as PILImage
                            pil = PILImage.open(io.BytesIO(img_bytes)).convert("RGB")
                            max_w = 1400
                            w, h = pil.size
                            if w > max_w:
                                pil.thumbnail((max_w, int(h * max_w / w)))
                            st.image(pil, use_container_width=True, caption=f"Slide {cur_slide.get('index') + 1}")
                            

                        except Exception:
                            try:
                                st.image(img_bytes, use_container_width=True, caption=f"Slide {cur_slide.get('index') + 1}")
                            except Exception:
                                st.info("Could not render image for this slide — showing extracted text instead.")
                                st.write(cur_slide.get("text", ""))
                    else:
                        st.info("No image available for this slide — showing extracted text")
                        st.write(cur_slide.get("text", ""))

                # Sidebar: slider + thumbnails (pagination for large decks)
                with side_col:
                    st.markdown("**Controls**")
                    if num_slides > 0:
                        s_val = st.slider("Go to slide", 1, num_slides, value=cur_idx + 1, key=slider_key)
                        new_idx = max(0, s_val - 1)
                        if new_idx != cur_idx:
                            st.session_state[viewer_index_key] = new_idx

                        st.markdown("---")

                        show_all = st.checkbox("Show all thumbnails (may be slow for large decks)", value=False)
                        per_page = st.number_input("Thumbnails per page", min_value=6, max_value=72, value=36, step=6, key=f"thumbs_pp_{uid}")
                        if show_all:
                            thumbs_to_show = slides_list
                        else:
                            # page-based thumbnails
                            page_key = f"thumbs_page_{uid}"
                            total_pages = max(1, math.ceil(num_slides / per_page))
                            page = st.number_input("Thumbnail page", 1, total_pages, value=1, key=page_key)
                            start = (page - 1) * per_page
                            thumbs_to_show = slides_list[start:start + per_page]

                        st.markdown("**Thumbnails**")
                        thumbs_per_row = 3
                        for row_start in range(0, len(thumbs_to_show), thumbs_per_row):
                            row = thumbs_to_show[row_start: row_start + thumbs_per_row]
                            cols = st.columns(len(row))
                            for ci, sthumb in enumerate(row):
                                t_idx = int(sthumb.get("index", row_start + ci))
                                t_img = (sthumb.get("images") or [None])[0]
                                with cols[ci]:
                                    if t_img:
                                        try:
                                            from PIL import Image as PILImage
                                            pil_t = PILImage.open(io.BytesIO(t_img)).convert("RGB")
                                            pil_t.thumbnail((260, 160))
                                            st.image(pil_t, use_container_width=True)
                                        except Exception:
                                            try:
                                                st.image(t_img, use_container_width=True)
                                            except Exception:
                                                st.markdown(f"<div class='small-muted'>Slide {t_idx+1}</div>", unsafe_allow_html=True)
                                    else:
                                        snippet = (sthumb.get("text") or "")[:100]
                                        st.markdown(
                                            f"<div class='small-muted' style='font-size:12px;padding:6px;border-radius:8px;border:1px solid rgba(255,255,255,0.02);'>{snippet}</div>",
                                            unsafe_allow_html=True,
                                        )
                                    open_key = f"open_thumb_{uid}_{t_idx}"
                                    if st.button("Open", key=open_key):
                                        st.session_state[viewer_index_key] = int(t_idx)
                    else:
                        st.info("No slides detected to display.")
                st.markdown("</div>", unsafe_allow_html=True)

        except Exception as e:
            st.error(f"Unexpected error during upload processing: {e}")
            log("Unhandled upload processing error:", e, getattr(e, "__traceback__", None))

# --------------------------
# --------------------------
# Lessons, Chat Q&A, Quizzes, Flashcards, Export, Progress, Settings
# Unified robust replacement block
# --------------------------
# --------------------------
# Lessons, Chat Q&A, Quizzes, Flashcards, Export, Progress, Settings
# Unified replacement — robust error handling, fallbacks, features
# --------------------------
import math

TOP_K = int(st.session_state.get("TOP_K", globals().get("TOP_K", 5)))
EMBED_MODEL_NAME = st.session_state.get("EMBEDDING_MODEL_NAME", globals().get("EMBEDDING_MODEL_NAME", "all-MiniLM-L6-v2"))

# Helper: safe get uploads list
uploads = st.session_state.get("uploads", []) or []

# --------------------------
# Lessons
# --------------------------
if active_tab == "Lessons":
    st.header("Generate Multi-level Lessons")
    uploads = st.session_state.get("uploads", []) or []
    if not uploads:
        st.info("No uploads yet. Go to Upload & Process.")
    else:
        # choose upload
        options = {u["id"]: u.get("filename", "untitled") for u in uploads}
        sel_id = st.selectbox("Select upload", options=list(options.keys()), format_func=lambda k: options[k], key="select_upload_lessons")
        upload = next((u for u in uploads if u["id"] == sel_id), None)
        if upload is None:
            st.warning("Selected upload not found.")
        else:
            st.markdown("**Context source for lesson generation**")
            context_mode = st.radio("", ["Indexed search (recommended)", "Whole document", "Specific slide"], index=0, key="lessons_context_mode")

            slides = upload.get("slides", []) or []
            max_idx = max([int(s.get("index", 0)) for s in slides]) if slides else 0

            seed_slide_idx = None
            if context_mode == "Specific slide":
                seed_slide_idx = st.number_input("Slide/Page index to focus on", min_value=0, max_value=max_idx, value=0, key="lessons_slide_index")
                slide_text = next((s.get("text", "") for s in slides if int(s.get("index", 0)) == int(seed_slide_idx)), "")
            elif context_mode == "Whole document":
                # join all slide texts
                slide_text = "\n\n".join([s.get("text", "") for s in slides])
            else:  # Indexed search
                # allow optional seed slide
                seed_with_slide = st.checkbox("Seed with a slide for relevance (optional)", value=False, key="lessons_seed_with_slide")
                if seed_with_slide:
                    seed_slide_idx = st.number_input("Seed slide index", min_value=0, max_value=max_idx, value=0, key="lessons_seed_slide_idx")
                    slide_text = next((s.get("text", "") for s in slides if int(s.get("index", 0)) == int(seed_slide_idx)), "")
                else:
                    # default seed -> whole document text
                    slide_text = "\n\n".join([s.get("text", "") for s in slides])

            st.markdown("<div class='card'>", unsafe_allow_html=True)
            st.subheader("Seed / preview text")
            st.write(slide_text if len(slide_text) < 4000 else slide_text[:4000] + "...")
            st.markdown("</div>", unsafe_allow_html=True)

            deep = st.checkbox("Produce a deeply detailed lesson (longer, step-by-step)", value=False, key="lessons_deep")
            auto_create = st.checkbox("Also auto-create quiz + flashcards and save to DB", value=True, key="lessons_autocreate")

            if st.button("Generate Lesson (Beginner → Advanced)", key="generate_lesson_btn"):
                # Build context snippets depending on mode and index availability
                snippets = []
                try:
                    if context_mode == "Indexed search (recommended)":
                        idx = upload.get("index")
                        if idx and upload.get("chunks"):
                            # compute embedding for seed text
                            try:
                                if "load_sentence_transformer" in globals() and callable(globals()["load_sentence_transformer"]):
                                    model = load_sentence_transformer(st.session_state.get("EMBEDDING_MODEL_NAME", EMBED_MODEL_NAME))
                                    q_emb = embed_texts(model, [slide_text])
                                    D, I = idx.search(q_emb, int(st.session_state.get("TOP_K", TOP_K)))
                                    for j in (I[0] if len(I) > 0 else []):
                                        if isinstance(j, int) and 0 <= j < len(upload.get("chunks", [])):
                                            slide_num = upload.get("mapping", [{}])[j].get("slide") if j < len(upload.get("mapping", [])) else None
                                            prefix = f"[Slide {slide_num}] " if slide_num is not None else ""
                                            snippets.append(prefix + upload["chunks"][j])
                            except Exception as e:
                                log("Indexed search failed (embedding/index):", e)
                        # fallback to whole doc if snippets empty (this allows AI even without index)
                        if not snippets:
                            snippets = [ "\n\n".join([s.get("text", "") for s in slides]) ]
                    elif context_mode == "Whole document":
                        snippets = ["\n\n".join([s.get("text", "") for s in slides])]
                    else:  # Specific slide
                        snippets = [slide_text]
                except Exception as e:
                    log("Building context snippets failed:", e)
                    snippets = ["\n\n".join([s.get("text", "") for s in slides])]

                related = "\n\n".join(snippets)

                # Call LLM via safe wrapper
                with st.spinner("Generating lesson..."):
                    try:
                        if deep:
                            lesson = _safe_call("generate_deep_lesson", slide_text, related, default=None)
                        else:
                            lesson = _safe_call("generate_multilevel_lesson", slide_text, related, default=None)
                        if lesson is None:
                            st.warning("Lesson generation helper not available in this runtime.")
                            lesson = "[Lesson generation not available]"
                    except Exception as e:
                        lesson = f"[Lesson generation failed: {e}]"
                        log("Lesson generation error:", e)

                st.subheader("Generated Lesson")
                st.markdown(lesson)

                # Auto-create artifacts
                if auto_create:
                    st.info("Auto-generating MCQs and flashcards from the generated lesson...")
                    mcqs = _safe_call("generate_mcq_set_from_text", lesson, qcount=8, default=[])
                    fcards = _safe_call("generate_flashcards_from_text", lesson, n=12, default=[])
                    cur = _get_db_cursor()
                    if cur is None:
                        st.warning("No DB available — generated artifacts will not be persisted.")
                    else:
                        try:
                            db_uid = upload_db_id(upload)
                            if db_uid is None:
                                st.warning("Unable to determine DB id for this upload — artifacts will not be saved.")
                            else:
                                for q in mcqs or []:
                                    try:
                                        opts = q.get("options", []) if isinstance(q.get("options", []), list) else []
                                        correct_index = int(q.get("answer_index", 0) if q.get("answer_index") is not None else 0)
                                        cur.execute(
                                            "INSERT INTO quizzes (upload_id, question, options, correct_index, created_at) VALUES (?, ?, ?, ?, ?)",
                                            (db_uid, q.get("question", ""), json.dumps(opts, ensure_ascii=False), correct_index, int(time.time())),
                                        )
                                    except Exception as ex_q:
                                        log("Failed to insert MCQ:", ex_q)
                                inserted = 0
                                for card in fcards or []:
                                    qtext = card.get("q") or card.get("question") or ""
                                    atext = card.get("a") or card.get("answer") or ""
                                    if qtext and atext:
                                        try:
                                            cur.execute(
                                                "INSERT INTO flashcards (upload_id, question, answer, easiness, interval, repetitions, next_review) VALUES (?, ?, ?, ?, ?, ?, ?)",
                                                (db_uid, qtext, atext, 2.5, 1, 0, int(time.time())),
                                            )
                                            inserted += 1
                                        except Exception as ex_f:
                                            log("Failed to insert flashcard:", ex_f)
                                _db_conn.commit()
                                st.success(f"Saved {len(mcqs or [])} MCQs and {inserted} flashcards to DB (if any).")
                        except Exception as e:
                            st.warning(f"Could not save generated artifacts: {e}")

                # Optional TTS export using gTTS if available
                if globals().get("_HAS_GTTS") and st.button("Export lesson as MP3 (TTS)", key="lessons_export_tts"):
                    try:
                        fname, mp3bytes = _safe_call("text_to_speech_download", lesson, default=(None, None))
                        if fname and mp3bytes:
                            st.download_button("Download lesson audio", mp3bytes, file_name=fname, mime="audio/mpeg", key=f"dl_lesson_tts")
                        else:
                            st.warning("TTS helper not available or returned no data.")
                    except Exception as e:
                        st.error(f"TTS failed: {e}")

# --------------------------
# Chat Q&A
# --------------------------
elif active_tab == "Chat Q&A":
    st.header("Ask questions about your upload (Retrieval + LLM)")
    uploads = st.session_state.get("uploads", []) or []
    if not uploads:
        st.info("No uploads yet. Upload files first.")
    else:
        options = {u["id"]: u.get("filename", "untitled") for u in uploads}
        sel_id = st.selectbox("Select upload", options=list(options.keys()), format_func=lambda k: options[k], key="select_upload_chat")
        upload = next((u for u in uploads if u["id"] == sel_id), None)
        if not upload:
            st.warning("Selected upload not found.")
        else:
            st.markdown("**Context source**")
            context_mode = st.radio("", ["Indexed search (recommended)", "Whole document", "Specific slide"], index=0, key="chat_context_mode")
            slides = upload.get("slides", []) or []
            max_idx = max([int(s.get("index", 0)) for s in slides]) if slides else 0

            seed_slide = None
            if context_mode == "Specific slide":
                seed_slide = st.number_input("Slide index for context", min_value=0, max_value=max_idx, value=0, key="chat_seed_slide")

            question = st.text_area("Ask a question about the slides/pages", key="chat_question")

            if st.button("Get Answer", key="chat_get_answer") and question.strip():
                top_ctx = []
                try:
                    if context_mode == "Indexed search (recommended)":
                        idx = upload.get("index")
                        if idx and upload.get("chunks"):
                            try:
                                model = load_sentence_transformer(st.session_state.get("EMBEDDING_MODEL_NAME", EMBED_MODEL_NAME))
                                q_emb = embed_texts(model, [question])
                                D, I = idx.search(q_emb, int(st.session_state.get("TOP_K", TOP_K)))
                                for j in (I[0] if len(I) > 0 else []):
                                    if isinstance(j, int) and 0 <= j < len(upload.get("chunks", [])):
                                        slide_num = upload.get("mapping", [{}])[j].get("slide") if j < len(upload.get("mapping", [])) else None
                                        prefix = f"[Slide {slide_num}] " if slide_num is not None else ""
                                        top_ctx.append(prefix + upload["chunks"][j])
                            except Exception as e:
                                log("Index Q&A failed:", e)
                        # fallback: whole document
                        if not top_ctx:
                            top_ctx = ["\n\n".join([s.get("text", "") for s in slides])]
                    elif context_mode == "Whole document":
                        top_ctx = ["\n\n".join([s.get("text", "") for s in slides])]
                    else:
                        top_ctx = [ next((s.get("text","") for s in slides if int(s.get("index", 0)) == int(seed_slide)), "") ]
                except Exception as e:
                    log("Building Q&A context failed:", e)
                    top_ctx = ["\n\n".join([s.get("text", "") for s in slides])]

                top_ctx_text = "\n\n".join(top_ctx)
                system = "You are a helpful tutor. Use the provided context to answer concisely; cite slide/page indices if possible."
                prompt = f"CONTEXT:\n{top_ctx_text}\n\nQUESTION:\n{question}\n\nAnswer concisely and provide one short example or analogy."

                with st.spinner("Querying LLM..."):
                    try:
                        ans = _safe_call("call_openrouter_chat", system, prompt, max_tokens=450, default="[OpenRouter call not available]")
                        st.subheader("Answer")
                        st.write(ans)
                    except Exception as e:
                        st.error(f"OpenRouter call failed: {e}")

                if top_ctx:
                    st.markdown("---")
                    st.markdown("**Context used (excerpt):**")
                    for j, s in enumerate(top_ctx[: int(st.session_state.get("TOP_K", TOP_K)) ]):
                        st.code(s[:800] + ("..." if len(s) > 800 else ""))

# --------------------------
# Quizzes
# --------------------------
elif active_tab == "Quizzes":
    st.header("Auto-generated Quizzes")
    uploads = st.session_state.get("uploads", []) or []
    if not uploads:
        st.info("No uploads yet.")
    else:
        options = {u["id"]: u.get("filename", "untitled") for u in uploads}
        sel_id = st.selectbox("Select upload", options=list(options.keys()), format_func=lambda k: options[k], key="select_upload_quizzes")
        upload = next((u for u in uploads if u["id"] == sel_id), None)
        if upload:
            slides = upload.get("slides", []) or []
            max_idx = max([int(s.get("index", 0)) for s in slides]) if slides else 0
            slide_idx = st.number_input("Slide/Page index to generate quiz from (or leave 0 for whole doc)", min_value=0, max_value=max_idx, value=0, key="quiz_slide_idx")
            mode = st.radio("Context mode", ["Slide (selected)", "Whole document"], index=1, key="quiz_context_mode")
            if mode == "Slide (selected)":
                context_text = next((s.get("text", "") for s in slides if int(s.get("index", 0)) == int(slide_idx)), "")
            else:
                context_text = "\n\n".join([s.get("text", "") for s in slides])

            if st.button("Generate Quiz (MCQs)", key="generate_quiz_btn"):
                with st.spinner("Creating MCQs..."):
                    qset = _safe_call("generate_mcq_set_from_text", context_text, qcount=5, default=[])
                if not qset:
                    st.warning("No MCQs generated (helper unavailable or empty result).")
                else:
                    st.success("Quiz ready — try it below")
                    # render each question + radio selection
                    user_answers = {}
                    for qi, q in enumerate(qset):
                        st.markdown(f"**Q{qi + 1}.** {q.get('question', '')}")
                        opts = q.get("options", [])
                        if not isinstance(opts, list) or len(opts) == 0:
                            opts = ["A", "B", "C", "D"]
                        choice = st.radio(f"Select Q{qi + 1}", opts, key=f"quiz_{sel_id}_{slide_idx}_{qi}")
                        user_answers[qi] = (choice, opts, int(q.get("answer_index", 0) if q.get("answer_index") is not None else 0))

                    # Submit all answers (one-shot)
                    if st.button("Submit all answers", key="submit_all_quiz"):
                        correct = 0
                        for qi, (choice, opts, correct_idx) in user_answers.items():
                            chosen_idx = opts.index(choice) if choice in opts else 0
                            if chosen_idx == correct_idx:
                                correct += 1
                        st.info(f"Score: {correct} / {len(user_answers)}")

                    # Persist quizzes to DB (best-effort)
                    cur = _get_db_cursor()
                    if cur:
                        try:
                            db_uid = upload_db_id(upload)
                            if db_uid:
                                for q in qset:
                                    cur.execute("INSERT INTO quizzes (upload_id, question, options, correct_index, created_at) VALUES (?, ?, ?, ?, ?)",
                                                (db_uid, q.get("question", ""), json.dumps(q.get("options", []), ensure_ascii=False), int(q.get("answer_index", 0) if q.get("answer_index") is not None else 0), int(time.time())))
                                _db_conn.commit()
                                st.success("Quiz saved to DB (if DB available).")
                        except Exception as e:
                            st.warning(f"Could not save quiz to DB: {e}")

# --------------------------
# Flashcards
# --------------------------
elif active_tab == "Flashcards":
    st.header("Flashcards & Spaced Repetition")
    uploads = st.session_state.get("uploads", []) or []
    if not uploads:
        st.info("No uploads yet. Go to Upload & Process.")
    else:
        options = {u["id"]: u.get("filename", "untitled") for u in uploads}
        sel_id = st.selectbox("Select upload to work with", options=list(options.keys()), format_func=lambda k: options[k], key="select_upload_flashcards")
        upload = next((u for u in uploads if u["id"] == sel_id), None)
        if not upload:
            st.warning("Selected upload not found.")
        else:
            slides = upload.get("slides", []) or []
            max_idx = max([int(s.get("index", 0)) for s in slides]) if slides else 0
            slide_idx = st.number_input("Slide/Page index to extract flashcards from", min_value=0, max_value=max_idx, value=0, key=f"flash_idx_{sel_id}")
            text = next((s.get("text", "") for s in slides if int(s.get("index", 0)) == int(slide_idx)), "")

            if st.button("Generate Flashcards from this slide/page", key=f"gen_flash_{sel_id}_{slide_idx}"):
                with st.spinner("Generating flashcards..."):
                    cards = _safe_call("generate_flashcards_from_text", text, n=12, default=[])
                if not cards:
                    st.warning("No flashcards generated (helper not available or returned no cards).")
                else:
                    cur = _get_db_cursor()
                    if cur is None:
                        st.warning("Database not available; generated flashcards will not be persisted.")
                    else:
                        inserted = 0
                        db_uid = upload_db_id(upload)
                        try:
                            for card in cards:
                                qtext = card.get("q") or card.get("question") or ""
                                atext = card.get("a") or card.get("answer") or ""
                                if qtext and atext:
                                    cur.execute(
                                        "INSERT INTO flashcards (upload_id, question, answer, easiness, interval, repetitions, next_review) VALUES (?, ?, ?, ?, ?, ?, ?)",
                                        (db_uid, qtext, atext, 2.5, 1, 0, int(time.time())),
                                    )
                                    inserted += 1
                            _db_conn.commit()
                            if inserted:
                                st.success(f"Inserted {inserted} flashcards into your deck.")
                                st.markdown("**Preview (first 5):**")
                                for i, card in enumerate(cards[:5]):
                                    qtxt = card.get("q") or card.get("question") or ""
                                    atxt = card.get("a") or card.get("answer") or ""
                                    st.markdown(f"**{i+1}.** {qtxt}")
                                    st.markdown(f"<div class='small-muted'>Answer: {atxt}</div>", unsafe_allow_html=True)
                            else:
                                st.info("No valid Q/A pairs were found in the generated output.")
                        except Exception as e:
                            st.error(f"Failed to save flashcards: {e}")
                            log("flashcard save error:", e)

            st.markdown("---")
            st.subheader("Review due flashcards")

            # Practice mechanism (single-card practice flow)
            cur = _get_db_cursor()
            db_uid = upload_db_id(upload) if cur else None
            now = int(time.time())
            due_cards = []
            if cur and db_uid is not None:
                try:
                    cur.execute(
                        "SELECT id, question, answer, easiness, interval, repetitions, next_review FROM flashcards WHERE upload_id = ? AND (next_review IS NULL OR next_review <= ?) ORDER BY next_review ASC",
                        (db_uid, now),
                    )
                    due_cards = cur.fetchall()
                except Exception as e:
                    log("Failed to fetch due cards:", e)
                    due_cards = []

            if not due_cards:
                st.info("No cards due for this upload. Generate some flashcards or wait for scheduled review.")
            else:
                sess_key = f"practice_idx_{sel_id}"
                st.session_state.setdefault(sess_key, 0)
                pidx = st.session_state[sess_key]
                pidx = max(0, min(pidx, max(0, len(due_cards) - 1)))
                st.session_state[sess_key] = pidx

                # normalize DB row -> card
                def _row_to_card(row):
                    try:
                        fid, qtext, atext, eas, inter, reps, nxt = row
                    except Exception:
                        fid = row[0]
                        qtext = row[1] if len(row) > 1 else ""
                        atext = row[2] if len(row) > 2 else ""
                        eas = row[3] if len(row) > 3 else 2.5
                        inter = row[4] if len(row) > 4 else 1
                        reps = row[5] if len(row) > 5 else 0
                        nxt = row[6] if len(row) > 6 else 0
                    return {"id": fid, "q": qtext, "a": atext, "eas": eas, "int": inter, "reps": reps, "next": nxt}

                card = _row_to_card(due_cards[pidx])
                st.markdown(f"**Card {pidx + 1} of {len(due_cards)}**")
                st.write(card["q"])

                # Show/Hide answer
                show_key = f"show_answer_{card['id']}"
                if st.button("Show Answer", key=show_key):
                    st.write(card["a"])

                rating = st.radio("How well did you recall? (0 = no recall, 5 = perfect)", options=[0, 1, 2, 3, 4, 5], index=3, key=f"rating_{card['id']}")
                if st.button("Submit Rating", key=f"submit_practice_{card['id']}"):
                    try:
                        eas_new, interval_new, reps_new, next_review = _safe_call("sm2_update_card", card["eas"], card["int"], card["reps"], int(rating), default=(card["eas"], card["int"], card["reps"], int(time.time())))
                        cur.execute("UPDATE flashcards SET easiness = ?, interval = ?, repetitions = ?, next_review = ? WHERE id = ?", (eas_new, interval_new, reps_new, next_review, card["id"]))
                        _db_conn.commit()
                        st.success("Card updated; next review scheduled.")
                        # advance to next card
                        st.session_state[sess_key] = min(pidx + 1, len(due_cards) - 1)
                    except Exception as e:
                        st.error(f"Failed to update card: {e}")

# --------------------------
# Export
# --------------------------
elif active_tab == "Export":
    st.header("Export — Anki / Audio / Raw")
    uploads = st.session_state.get("uploads", []) or []
    if not uploads:
        st.info("No uploads available.")
    else:
        options = {u["id"]: u.get("filename", "untitled") for u in uploads}
        sel_id = st.selectbox("Select upload to export from", options=list(options.keys()), format_func=lambda k: options[k], key="select_upload_export")
        upload = next((u for u in uploads if u["id"] == sel_id), None)
        if not upload:
            st.warning("Selected upload not found.")
        else:
            st.markdown("**Flashcards / Anki**")
            if st.button("Export flashcards to Anki (TSV)", key=f"export_anki_{sel_id}"):
                try:
                    res = _safe_call("anki_export_csv_for_upload", upload_db_id(upload), _db_conn, default=None)
                    if res:
                        fname, data = res
                        st.download_button("Download Anki TSV", data, file_name=fname, mime="text/tab-separated-values", key=f"dl_anki_{sel_id}")
                    else:
                        st.warning("anki_export_csv_for_upload helper not available.")
                except Exception as e:
                    st.error(f"Export failed: {e}")

            st.markdown("---")
            if globals().get("_HAS_GTTS"):
                if st.button("Export all generated lessons as MP3 (single)", key=f"export_mp3_{sel_id}"):
                    lesson_text = ""
                    for s in upload.get("slides", [])[:1000]:
                        lesson_text += f"Slide {s.get('index', '?')}. {s.get('text','')}\n\n"
                    try:
                        res = _safe_call("text_to_speech_download", lesson_text, default=None)
                        if res:
                            fname, data = res
                            st.download_button("Download lessons MP3", data, file_name=fname, mime="audio/mpeg", key=f"dl_mp3_{sel_id}")
                        else:
                            st.warning("TTS helper not available.")
                    except Exception as e:
                        st.error(f"TTS failed: {e}")
            else:
                st.info("gTTS not available — TTS export disabled.")

            st.markdown("---")
            if st.button("Download extracted slides (JSON)", key=f"export_json_{sel_id}"):
                try:
                    ark = json.dumps(upload.get("slides", []), ensure_ascii=False, indent=2).encode("utf-8")
                    st.download_button("Download JSON", ark, file_name=f"{upload['filename']}_extracted.json", mime="application/json", key=f"dl_json_{sel_id}")
                except Exception as e:
                    st.error(f"Could not prepare download: {e}")

# --------------------------
# Progress
# --------------------------
elif active_tab == "Progress":
    st.header("Progress & Analytics")
    st.markdown("Overview of uploads and study artifacts")

    cur = _get_db_cursor()
    rows = []
    if cur is not None:
        try:
            cur.execute("SELECT id, filename, uploaded_at, meta FROM uploads ORDER BY uploaded_at DESC")
            rows = cur.fetchall()
        except Exception as e:
            log("Failed to fetch uploads from DB:", e)
            rows = []

    if not rows:
        st.info("No uploads logged in DB yet.")
    else:
        for r in rows:
            try:
                uid, fname, uploaded_at, meta = r
            except Exception:
                uid = r[0]
                fname = r[1] if len(r) > 1 else "unknown"
                uploaded_at = r[2] if len(r) > 2 else 0
                meta = r[3] if len(r) > 3 else "{}"
            st.markdown(f"**{fname}** — uploaded at {time.ctime(int(uploaded_at)) if uploaded_at else 'unknown'}")
            try:
                meta_obj = json.loads(meta) if isinstance(meta, str) else (meta or {})
            except Exception:
                meta_obj = {}
            st.write(meta_obj)
            c2 = _get_db_cursor()
            if c2:
                try:
                    c2.execute("SELECT COUNT(*) FROM flashcards WHERE upload_id = ?", (uid,))
                    fc_count = int(c2.fetchone()[0] or 0)
                    c2.execute("SELECT COUNT(*) FROM quizzes WHERE upload_id = ?", (uid,))
                    q_count = int(c2.fetchone()[0] or 0)
                except Exception as e:
                    log("Failed counting artifacts:", e)
                    fc_count = 0
                    q_count = 0
            else:
                fc_count = 0
                q_count = 0
            st.write(f"Flashcards: {fc_count} • Quizzes: {q_count}")
            st.markdown("---")

    c = _get_db_cursor()
    total_fc = total_q = 0
    if c:
        try:
            c.execute("SELECT COUNT(*) FROM flashcards")
            total_fc = int(c.fetchone()[0] or 0)
            c.execute("SELECT COUNT(*) FROM quizzes")
            total_q = int(c.fetchone()[0] or 0)
        except Exception as e:
            log("Failed to fetch totals:", e)
    st.write(f"Total flashcards in DB: {total_fc}")
    st.write(f"Total quizzes in DB: {total_q}")

# --------------------------
# Settings
# --------------------------
elif active_tab == "Settings":
    st.header("Settings & Diagnostics")

    # Ensure session API key key exists
    st.session_state.setdefault("OPENROUTER_API_KEY", os.getenv("OPENROUTER_API_KEY", globals().get("DEFAULT_OPENROUTER_KEY", "")))
    key = st.text_input("OpenRouter API key (leave blank to use env/default)", value=st.session_state.get("OPENROUTER_API_KEY", ""), type="password", key="api_key_input")
    if st.button("Save API Key (session only)", key="save_api_key"):
        st.session_state["OPENROUTER_API_KEY"] = (key.strip() or st.session_state.get("OPENROUTER_API_KEY", ""))
        st.success("API key set for this session.")

    # Model & Index settings
    st.markdown("Model & Index settings")
    st.session_state.setdefault("EMBEDDING_MODEL_NAME", globals().get("EMBEDDING_MODEL_NAME", "all-MiniLM-L6-v2"))
    st.session_state.setdefault("TOP_K", globals().get("TOP_K", 5))

    emb_name = st.text_input("Embedding model name", value=st.session_state.get("EMBEDDING_MODEL_NAME"), key="setting_emb_model")
    topk_val = st.number_input("Search TOP_K", min_value=1, max_value=50, value=int(st.session_state.get("TOP_K", 5)), key="setting_top_k")

    if st.button("Save Settings", key="save_settings"):
        st.session_state["EMBEDDING_MODEL_NAME"] = emb_name.strip() or st.session_state["EMBEDDING_MODEL_NAME"]
        st.session_state["TOP_K"] = int(topk_val)
        st.success("Settings saved to session.")

    # Diagnostics
    st.markdown("Diagnostics:")
    st.write({
        "faiss_available": bool(globals().get("_HAS_FAISS", False)),
        "pymupdf_available": bool(globals().get("_HAS_PYMUPDF", False)),
        "easyocr_available": bool(globals().get("_HAS_EASYOCR", False)),
        "sentence_transformers_available": bool(globals().get("_HAS_SENTENCE_TRANSFORMERS", False)),
        "gtts_available": bool(globals().get("_HAS_GTTS", False))
    })

    # OpenRouter test
    if st.button("Test OpenRouter (small ping)", key="test_openrouter"):
        try:
            ping = _safe_call("call_openrouter_chat", "You are a test bot.", "Say 'pong' in a plain short reply.", max_tokens=20, default="[OpenRouter not available]")
            st.code(ping)
        except Exception as e:
            st.error(f"OpenRouter test failed: {e}")

    if st.button("Clear all session uploads (session state only)", key="clear_uploads_btn"):
        st.session_state["uploads"] = []
        st.success("Cleared session uploads.")
